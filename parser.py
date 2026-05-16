"""Parse uploaded resources into reusable text chunks."""

from dataclasses import dataclass
import os
import re
from typing import List

import pdfplumber
from pptx import Presentation


@dataclass
class RawChunk:
    text: str
    source_file: str
    chunk_index: int
    topic_hint: str = ""
    page_or_slide: int = 0


UNIT_HEADING_RE = re.compile(
    r"^\s*(unit\s*[-–]?\s*(i{1,3}|iv|v?i{0,3}|[1-5]))\b.*$",
    re.IGNORECASE | re.MULTILINE,
)


def _basename(filepath: str) -> str:
    return os.path.basename(filepath)


def detect_source_type(filepath: str) -> str:
    """Classify uploads as syllabus, textbook, or ppt based on file contents."""
    ext = os.path.splitext(filepath)[1].lower().lstrip(".")
    if ext in {"ppt", "pptx"}:
        return "ppt"

    full_text = "\n".join(page.extract_text() or "" for page in _open_pdf_pages(filepath))
    return "syllabus" if len(UNIT_HEADING_RE.findall(full_text)) >= 2 else "textbook"


def _open_pdf_pages(filepath: str):
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            yield page


# ── PPT/PPTX ──────────────────────────────────────────────────────────────────

def parse_pptx(filepath: str) -> List[RawChunk]:
    prs = Presentation(filepath)
    chunks = []
    filename = _basename(filepath)

    for i, slide in enumerate(prs.slides):
        texts = []
        title = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if not line:
                    continue
                # first non-empty text on the slide is treated as the title
                if not title:
                    title = line
                texts.append(line)

        combined = "\n".join(texts).strip()
        if not combined:
            continue

        chunks.append(RawChunk(
            text=combined,
            source_file=filename,
            chunk_index=i,
            topic_hint=title[:120],
            page_or_slide=i + 1,
        ))

    return chunks


# ── PDF ───────────────────────────────────────────────────────────────────────

def _split_by_units(full_text: str, filename: str) -> List[RawChunk]:
    """Split syllabus PDFs on unit headings while keeping the heading as context."""
    matches = list(UNIT_HEADING_RE.finditer(full_text))
    if not matches:
        return []

    chunks = []
    for idx, match in enumerate(matches):
        start = match.start()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(full_text)
        body = full_text[start:end].strip()
        if not body:
            continue
        chunks.append(
            RawChunk(
                text=body,
                source_file=filename,
                chunk_index=idx,
                topic_hint=match.group(0).strip()[:120],
                page_or_slide=0,
            )
        )
    return chunks


def _fixed_window_chunks(
    pages: List[str],
    filename: str,
    chunk_size: int = 512,
    overlap: int = 100,
) -> List[RawChunk]:
    """Word-based chunking for textbook PDFs when unit headings are absent."""
    words: List[str] = []
    word_page_map: List[int] = []

    for page_num, page_text in enumerate(pages, start=1):
        page_words = re.findall(r"\S+", page_text or "")
        words.extend(page_words)
        word_page_map.extend([page_num] * len(page_words))

    if not words:
        return []

    chunks = []
    start = 0
    chunk_idx = 0
    step = max(1, chunk_size - min(overlap, chunk_size - 1))

    while start < len(words):
        end = min(start + chunk_size, len(words))
        text = " ".join(words[start:end]).strip()
        if text:
            chunks.append(
                RawChunk(
                    text=text,
                    source_file=filename,
                    chunk_index=chunk_idx,
                    topic_hint="",
                    page_or_slide=word_page_map[start],
                )
            )
            chunk_idx += 1
        if end >= len(words):
            break
        start += step

    return chunks


def parse_pdf(filepath: str, chunk_size: int = 512, overlap: int = 100) -> List[RawChunk]:
    filename = _basename(filepath)
    pages_text: List[str] = []

    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages_text.append(text)

    full_text = "\n".join(pages_text)

    if len(UNIT_HEADING_RE.findall(full_text)) >= 2:
        return _split_by_units(full_text, filename)
    else:
        return _fixed_window_chunks(pages_text, filename, chunk_size, overlap)


# ── Dispatcher ────────────────────────────────────────────────────────────────

def parse_file(filepath: str, chunk_size: int = 512, overlap: int = 100) -> List[RawChunk]:
    ext = os.path.splitext(filepath)[1].lower().lstrip(".")
    if ext in ("ppt", "pptx"):
        return parse_pptx(filepath)
    elif ext == "pdf":
        return parse_pdf(filepath, chunk_size, overlap)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")